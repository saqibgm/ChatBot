"""
Professional PowerPoint Presentations - Clean Corporate Colors
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ============== PROFESSIONAL CORPORATE COLORS ==============

# GLPI Project - Clean Navy Blue Theme
GLPI_COLORS = {
    'bg': RGBColor(0xff, 0xff, 0xff),        # White background
    'header_bg': RGBColor(0x1e, 0x3a, 0x5f),  # Navy blue
    'accent': RGBColor(0x1e, 0x3a, 0x5f),     # Navy blue
    'text': RGBColor(0x2c, 0x3e, 0x50),       # Dark gray-blue
    'table_header': RGBColor(0x1e, 0x3a, 0x5f),  # Navy
    'table_bg': RGBColor(0xf8, 0xf9, 0xfa),   # Light gray
    'table_alt': RGBColor(0xe9, 0xec, 0xef),  # Alternate row
    'white': RGBColor(0xff, 0xff, 0xff),
    'light_text': RGBColor(0x6c, 0x75, 0x7d),  # Gray
}

# AI Coding - Clean Blue-Gray Theme
AI_COLORS = {
    'bg': RGBColor(0xff, 0xff, 0xff),         # White background
    'header_bg': RGBColor(0x2c, 0x3e, 0x50),   # Dark blue-gray
    'accent': RGBColor(0x34, 0x98, 0xdb),      # Bright blue
    'text': RGBColor(0x2c, 0x3e, 0x50),        # Dark gray
    'table_header': RGBColor(0x34, 0x98, 0xdb),  # Blue
    'table_bg': RGBColor(0xf8, 0xf9, 0xfa),    # Light gray
    'table_alt': RGBColor(0xe9, 0xec, 0xef),   # Alternate row
    'white': RGBColor(0xff, 0xff, 0xff),
    'light_text': RGBColor(0x6c, 0x75, 0x7d),  # Gray
}


def add_title_slide(prs, title, subtitle, colors):
    """Clean professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Navy/Blue header band
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(3.5))
    header.fill.solid()
    header.fill.fore_color.rgb = colors['header_bg']
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    tf = sub_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(subtitle.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = colors['text']
        p.alignment = PP_ALIGN.CENTER
    
    # Bottom accent line
    line = slide.shapes.add_shape(1, Inches(3), Inches(3.5), Inches(4), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = colors['accent']
    line.line.fill.background()
    
    return slide


def add_content_slide(prs, title, content_lines, colors, emoji=""):
    """Clean content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = colors['header_bg']
    header.line.fill.background()
    
    # Title
    title_text = f"{emoji}  {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if line.startswith("   "):
            p.font.size = Pt(22)
            p.font.color.rgb = colors['light_text']
        elif line == "":
            p.font.size = Pt(12)
        else:
            p.font.size = Pt(26)
            p.font.color.rgb = colors['text']
        p.space_after = Pt(8)
    
    return slide


def add_two_column_slide(prs, title, left_content, right_content, colors, emoji=""):
    """Two-column layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = colors['header_bg']
    header.line.fill.background()
    
    # Title
    title_text = f"{emoji}  {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # Left box background
    left_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(1.5), Inches(4.4), Inches(5.5))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = colors['table_bg']
    left_bg.line.fill.background()
    
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(4), Inches(5.1))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if line.startswith("•"):
            p.font.size = Pt(22)
            p.font.color.rgb = colors['light_text']
        else:
            p.font.size = Pt(24)
            p.font.color.rgb = colors['accent']
            p.font.bold = True
        p.space_after = Pt(6)
    
    # Right box background
    right_bg = slide.shapes.add_shape(1, Inches(5.2), Inches(1.5), Inches(4.4), Inches(5.5))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = colors['table_bg']
    right_bg.line.fill.background()
    
    right_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4), Inches(5.1))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if line.startswith("•"):
            p.font.size = Pt(22)
            p.font.color.rgb = colors['light_text']
        else:
            p.font.size = Pt(24)
            p.font.color.rgb = colors['accent']
            p.font.bold = True
        p.space_after = Pt(6)
    
    return slide


def add_table_slide(prs, title, headers, rows, colors, emoji=""):
    """Professional table slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = colors['header_bg']
    header.line.fill.background()
    
    # Title
    title_text = f"{emoji}  {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = colors['white']
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    row_height = min(0.65, 5.2 / num_rows)
    
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.4), Inches(1.5),
        Inches(9.2), Inches(row_height * num_rows)
    ).table
    
    for i in range(num_cols):
        table.columns[i].width = int(Inches(9.2) / num_cols)
    
    # Header row
    for i, hdr in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors['table_header']
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(18)
            para.font.color.rgb = colors['white']
            para.alignment = PP_ALIGN.CENTER
    
    # Data rows with alternating colors
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors['table_bg'] if row_idx % 2 == 0 else colors['table_alt']
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(16)
                para.font.color.rgb = colors['text']
                para.alignment = PP_ALIGN.CENTER
    
    return slide


def create_glpi_presentation():
    """Project IQ Chatbot - Updated Features"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    C = GLPI_COLORS
    
    add_title_slide(prs, "Project IQ Chatbot",
        "Intelligent IT Support Automation\n\nRasa 3.6 • React • PostgreSQL", C)
    
    add_table_slide(prs, "Technology Stack",
        ["Component", "Technology"],
        [
            ["NLU Engine", "Rasa 3.6 + DIETClassifier"],
            ["Frontend", "React 19 + Tailwind CSS"],
            ["Backend", "Flask REST API"],
            ["Database", "PostgreSQL 18"],
            ["Integration", "GLPI REST API"],
        ], C, "🔧")
    
    add_content_slide(prs, "Core Features", [
        "✅  Smart Form Inference",
        "   Auto-detects Hardware/Network issues from text",
        "",
        "✅  Single-Message Creation",
        "   \"Urgent, my printer is jammed\" (One-shot)",
        "",
        "✅  Ticket Lifecycle Management",
        "   Create • Check Status • Update • List • Attach Files",
        "",
        "✅  Modern Chat Interface",
        "   Real-time • Markdown • Charts",
    ], C, "🎯")
    
    add_two_column_slide(prs, "Advanced Features",
        ["Chat Interface", "• Glassmorphism UI", "• Interactive Charts", "• File Drag & Drop", "• YouTube Embeds"],
        ["Analytics Dashboard", "• Conversation Stats", "• Intent Distribution", "• Configuration Panel"], C, "✨")
    
    add_table_slide(prs, "Project Metrics",
        ["Metric", "Value"],
        [["Intents", "28"], ["Slots", "16"], ["Custom Actions", "18+"], ["Lines of Code", "5,200+"]], C, "📊")
    
    add_table_slide(prs, "AI Development Speed",
        ["Task", "Before", "After", "Gain"],
        [["API Wrapper", "3 days", "4 hrs", "8×"], ["Chat UI", "1 week", "1 day", "5×"], ["Bug Fixing", "2 days", "1 hr", "16×"]], C, "⚡")
    
    add_table_slide(prs, "Challenges Solved",
        ["Problem", "Solution"],
        [["Rule Conflicts", "Cleaned Stories + Relaxed Policy"], ["Metadata Errors", "Defensive Coding in Action Server"], ["Confirmation", "Specific Affirm Rules"]], C, "🔨")
    
    add_content_slide(prs, "Live Demo Flow", [
        "1.  🔐  Login (Stable & Verified)",
        "",
        "2.  📝  Smart Ticket Creation",
        "   \"Urgent, printer jammed\" -> KB Check -> Form -> Confirm",
        "",
        "3.  🛑  Cancellation Capability",
        "   Stop at any point with \"Cancel\"",
        "",
        "4.  📊  Check Status & List Tickets",
    ], C, "🎬")
    
    add_two_column_slide(prs, "Roadmap 2025",
        ["Q1 - Generative & Proactive", "• RAG Knowledge Base", "• Proactive Systems Alerts", "• Human Agent Handoff"],
        ["Q2 - Omnichannel & Voice", "• MS Teams Integration", "• Voice Input (Whisper)", "• Sentiment Analysis"], C, "🗺️")
    
    add_content_slide(prs, "UI Highlights", [
        "🎨  Glassmorphism - Modern translucent design",
        "",
        "📊  Charts - Bar, Pie, Line visualization",
        "",
        "📎  Drag & Drop - Easy file uploads",
        "",
        "🌙  Dark Mode - Eye-friendly theme",
    ], C, "🖼️")
    
    add_table_slide(prs, "vs Competitors",
        ["Feature", "Project IQ", "Zendesk", "ServiceNow"],
        [["GLPI Native", "✅", "❌", "❌"], ["Self-Hosted", "✅", "❌", "⚠️"], ["Cost/Month", "$50", "$150+", "$500+"]], C, "⚔️")
    
    add_table_slide(prs, "Return on Investment",
        ["Metric", "Traditional", "AI-Assisted", "Savings"],
        [["Dev Hours", "400 hrs", "160 hrs", "60%"], ["Time to MVP", "8 weeks", "3 weeks", "63%"]], C, "💰")
    
    add_content_slide(prs, "Key Results", [
        "⚡  Stable, Reliable Core Flows",
        "",
        "📈  Zero Fallback on Commands",
        "",
        "🎯  Intelligent Intent & Entity Extraction",
        "",
        "🚀  Ready for Deployment",
    ], C, "🏆")
    
    # Next To Dos - Feature Roadmap with AI-Assisted Estimates
    add_table_slide(prs, "Next To Dos",
        ["Feature", "AI Time", "Your Time", "Priority"],
        [
            ["RAG Knowledge Base", "~2 hrs", "~1 hr", "🔴 Critical"],
            ["Human Agent Handoff", "~1.5 hrs", "~30 min", "🔴 Critical"],
            ["MS Teams / Slack", "~35 min", "~30 min", "🟡 High"],
            ["Proactive Alerts", "~1.5 hrs", "~30 min", "🟡 High"],
            ["Voice Input (STT)", "~1 hr", "~15 min", "🟢 Medium"],
            ["Image OCR/Vision", "~1.5 hrs", "~30 min", "🟢 Medium"],
            ["Sentiment Analysis", "~1 hr", "~15 min", "🟢 Medium"],
        ], C, "📋")
    
    add_title_slide(prs, "Questions?",
        "Project IQ Chatbot Demo\n\nThank you!", C)
    
    path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Project_IQ _Chatbot.pptx")
    prs.save(path)
    print(f"✅ Project IQ: {path}")


def create_ai_coding_presentation():
    """AI Coding - 18 slides with Blue-Gray theme - Complete Tools List"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    C = AI_COLORS
    
    add_title_slide(prs, "AI-Assisted Coding",
        "Tools, Techniques & Best Practices\n\nDeveloper Guide 2024", C)
    
    add_content_slide(prs, "Why AI Coding?", [
        "🚀  5-10× faster development",
        "",
        "🎯  Consistent quality & patterns",
        "",
        "📚  Learn new frameworks quickly",
        "",
        "🐛  Debug issues in minutes",
    ], C, "🤖")
    
    # AI IDEs & Editors
    add_table_slide(prs, "AI-Powered IDEs & Editors",
        ["Tool", "Best For", "Key Strength"],
        [
            ["Cursor", "Full AI development", "Multi-file editing"],
            ["Windsurf (Codeium)", "Fast AI coding", "Flow feature"],
            ["Antigravity (Google)", "Agentic development", "Full project context"],
            ["GitHub Copilot", "Code autocomplete", "VS Code/JetBrains"],
            ["Replit Agent", "Full app generation", "Browser-based"],
        ], C, "💻")
    
    # AI Chatbots
    add_table_slide(prs, "AI Chatbots & Assistants",
        ["Tool", "Best For", "Key Strength"],
        [
            ["Claude (Anthropic)", "Architecture, debugging", "200K context"],
            ["ChatGPT / GPT-4o", "General coding", "Code interpreter"],
            ["Google Gemini", "Multimodal, research", "1M+ context, free"],
            ["Microsoft Copilot", "Office + coding", "Free, Bing search"],
            ["Perplexity", "Documentation lookup", "Real-time web"],
        ], C, "💬")
    
    # App Builders
    add_table_slide(prs, "AI App Builders (No/Low Code)",
        ["Tool", "Best For", "Key Strength"],
        [
            ["v0.dev (Vercel)", "UI components", "React/Next.js"],
            ["Bolt.new", "Full-stack apps", "Instant preview"],
            ["Lovable", "Complete apps", "End-to-end dev"],
            ["Devin (Cognition)", "Autonomous dev", "Plans & codes alone"],
        ], C, "🏗️")
    
    # Specialized Tools
    add_table_slide(prs, "Specialized AI Tools",
        ["Tool", "Purpose"],
        [
            ["Amazon CodeWhisperer", "AWS dev + security scan"],
            ["Tabnine", "Enterprise, on-premise"],
            ["Cody (Sourcegraph)", "Large codebase search"],
            ["Codium AI", "Test generation"],
            ["Aider", "Git-aware terminal AI"],
            ["Blackbox AI", "Code search & generation"],
        ], C, "🔧")
    
    # Which tool guide
    add_content_slide(prs, "Which Tool for Which Task?", [
        "📝  Writing Code → Cursor, Copilot, Windsurf",
        "",
        "🏗️  Full App Generation → Replit, Bolt.new, v0",
        "",
        "🔍  Debugging → Claude, ChatGPT, Gemini",
        "",
        "📚  Learning → ChatGPT, Gemini, Phind",
        "",
        "🎨  UI Generation → v0.dev, Bolt.new",
    ], C, "🎯")
    
    add_content_slide(prs, "How AI Helps", [
        "⚡  Rapid Prototyping",
        "   Scaffold projects in minutes",
        "",
        "🔧  Complex Problems",
        "   Debug integrations faster",
        "",
        "📋  Code Quality",
        "   Consistent patterns throughout",
    ], C, "💪")
    
    add_table_slide(prs, "AI for Existing Code",
        ["Task", "AI Helps With", "Savings"],
        [["Understanding", "Explain functions", "70%"], ["Refactoring", "Modernize code", "60%"],
         ["Bug Fixing", "Trace issues", "80%"], ["Documentation", "Generate docs", "90%"]], C, "📂")
    
    add_content_slide(prs, "Effective Prompts", [
        "✅  GOOD PROMPT:",
        '   "Fix the user_id slot mapping in Rasa 3.6"',
        "",
        "❌  BAD PROMPT:",
        '   "Fix my code"',
        "",
        "💡  Include: Context, versions, examples",
    ], C, "💬")
    
    add_table_slide(prs, "Common Issues & Fixes",
        ["Problem", "Solution"],
        [["Hallucinated APIs", "Verify in docs"], ["Outdated Code", "Specify versions"],
         ["Security Gaps", "Always review"], ["Over-Engineering", "Ask for minimal"]], C, "⚠️")
    
    add_content_slide(prs, "Best Practices - DO", [
        "✅  Be specific with prompts",
        "",
        "✅  Provide context always",
        "",
        "✅  Iterate in small steps",
        "",
        "✅  Review all generated code",
    ], C, "👍")
    
    add_content_slide(prs, "Best Practices - DON'T", [
        "❌  Blindly copy-paste",
        "",
        "❌  Skip testing",
        "",
        "❌  Over-rely on AI",
        "",
        "❌  Share API keys or secrets",
    ], C, "👎")
    
    add_table_slide(prs, "Security Considerations",
        ["Concern", "Mitigation"],
        [["Code Privacy", "Enterprise versions"], ["API Keys", "Never share"],
         ["IP Ownership", "Check ToS"], ["Compliance", "SOC2 vendors"]], C, "🔒")
    
    add_two_column_slide(prs, "Team Impact",
        ["New Skills Needed", "• Prompt engineering", "• AI output review", "• Tool evaluation"],
        ["Productivity Gains", "• 40% faster reviews", "• Days not weeks", "• 24/7 pair programmer"], C, "👥")
    
    add_table_slide(prs, "Typical ROI",
        ["Metric", "Before", "After", "Gain"],
        [["Dev Hours", "400", "160", "60%"], ["MVP Time", "8 wks", "3 wks", "63%"]], C, "💰")
    
    add_content_slide(prs, "Key Takeaways", [
        "1️⃣  AI accelerates, doesn't replace",
        "",
        "2️⃣  Clear prompts = better output",
        "",
        "3️⃣  Always review and test",
        "",
        "4️⃣  Keep learning - tools evolve fast",
    ], C, "🎯")
    
    add_title_slide(prs, "Questions?",
        "Start small • Iterate fast • Always review\n\nThank you!", C)
    
    path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "AI_Coding_Complete.pptx")
    prs.save(path)
    print(f"✅ AI Coding: {path}")


if __name__ == "__main__":
    create_glpi_presentation()
    create_ai_coding_presentation()
    print("\n✨ Both presentations created!")
