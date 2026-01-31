"""
Script to update Createl_Chatbot_PPT.pptx:
1. Rename "Createl IT Support Chatbot" to "Createl - Service and Support"
2. Copy slides from AI_Coding_Complete.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy
import os

# File paths
createl_ppt_path = r"d:\99T\Rasa\my-createl-bot\Createl_Chatbot_PPT.pptx"
ai_ppt_path = r"d:\99T\Rasa\my-createl-bot\AI_Coding_Complete.pptx"
output_path = r"d:\99T\Rasa\my-createl-bot\Createl_Chatbot_PPT_Updated.pptx"

# Load presentations
print("Loading presentations...")
createl_prs = Presentation(createl_ppt_path)
ai_prs = Presentation(ai_ppt_path)

# Step 1: Find and replace title text
print("Renaming title from 'Createl IT Support Chatbot' to 'Createl - Service and Support'...")
old_title = "Createl IT Support Chatbot"
new_title = "Createl - Service and Support"

replacements_made = 0
for slide in createl_prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_title in run.text:
                        run.text = run.text.replace(old_title, new_title)
                        replacements_made += 1
                        print(f"  Replaced text in slide")

print(f"Total replacements: {replacements_made}")

# Step 2: Find specific slides in AI presentation
print("\nLooking for slides to copy from AI_Coding_Complete.pptx...")
slides_to_copy = ["How AI Helps", "Team Impact", "Key Takeaways"]
found_slides = []

for idx, slide in enumerate(ai_prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            for target in slides_to_copy:
                if target.lower() in text.lower():
                    found_slides.append((idx, target, slide))
                    print(f"  Found slide '{target}' at index {idx}")
                    break

print(f"Found {len(found_slides)} slides to copy")

# Step 3: Copy slides (python-pptx doesn't have direct slide copy, so we'll note the slide numbers)
# Actually, copying slides between presentations is complex in python-pptx
# We'll create a note about which slides to copy manually

print("\n--- SLIDE COPY INFORMATION ---")
print("python-pptx has limited support for copying slides between presentations.")
print("The following slides from AI_Coding_Complete.pptx should be copied:")
for idx, title, slide in found_slides:
    print(f"  Slide {idx + 1}: {title}")

# Save the updated Createl presentation (with title change)
print(f"\nSaving updated presentation to: {output_path}")
createl_prs.save(output_path)
print("Done! Title has been updated.")
print(f"\nTo complete the slide copy, please manually copy slides {[x[0]+1 for x in found_slides]} from AI_Coding_Complete.pptx")
